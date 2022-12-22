# -*- coding: UTF-8 -*-
from pptx import Presentation
import matplotlib.pyplot as plt
from pptx.util import Inches 


#開啟新的簡報物件
prs = Presentation()
#建立簡報檔第一張頁面物件
title_slide_layout = prs.slide_layouts[0] 
#增加一張簡報
slide = prs.slides.add_slide(title_slide_layout)
#設定第一張簡報的標題 
title = slide.shapes.title
title.text = "Hello Python PPT"
#設定第一張簡報的副標題
subtitle = slide.placeholders[1]
subtitle.text = "作者：Meiko 2020/10/01"
#將簡報物件存檔



def add_slide(prs, layout, img):
    #加一張投影片
	slide = prs.slides.add_slide(layout)
	shapes = slide.shapes
	
    #投影片標題
	title_shape = shapes.title
	title_shape.text = '長方圖示範'


	# show the figure
	left = Inches(3)
	height = Inches(4.5)
	left = top = Inches(3)
	pic = slide.shapes.add_picture(img, left, top, height=height)

	return slide


x_labels = ['A','B','C']
sales_num = [120,90,60]
plt.bar(x_labels,sales_num)
plt.savefig('graph.jpg')
img = 'graph.jpg'


title_slide_layout = prs.slide_layouts[1]
slide = add_slide(prs, title_slide_layout ,img)

prs.save("python_ppt_v1.pptx")